"""CV Final 발표 슬라이드 자동 생성 v5 (13장, U-Net 강화 중심).

박한샘 교수 평가 기준 + U-Net 강화 핵심 framing:
- 슬라이드 줄임 (16 → 13장)
- 이미지 최소화 (샘플 + 결과 1쌍씩)
- 모델 구조 설명 강화 (슬라이드 4)
- U-Net 강화 방법론 슬라이드 NEW (슬라이드 5)
- 디스커션 슬라이드 NEW (슬라이드 11)

출력: 바탕화면\\발표자료_CV_Final\\발표슬라이드_v5_최종.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

# ─── 경로 ───
SAMPLES_DIR = Path(__file__).resolve().parents[2] / "samples"
OUTPUT_DIR = Path("C:/Users/User/Desktop/발표자료_CV_Final")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
PPTX_PATH = OUTPUT_DIR / "발표슬라이드_v5_최종.pptx"

# ─── 색상 ───
NAVY = RGBColor(0x1E, 0x27, 0x61)
ACCENT = RGBColor(0xF9, 0x61, 0x67)
CHARCOAL = RGBColor(0x2C, 0x3E, 0x50)
MUTED = RGBColor(0x7B, 0x8F, 0xA1)
LIGHT_BG = RGBColor(0xFA, 0xFB, 0xFD)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x48, 0xBB, 0x78)
GOLD = RGBColor(0xF5, 0xA6, 0x23)

FONT_BOLD = "맑은 고딕"
FONT_REGULAR = "맑은 고딕"

SLIDE_W = Cm(33.867)
SLIDE_H = Cm(19.05)
TOTAL = 13


def make_pres():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=CHARCOAL,
             align=PP_ALIGN.LEFT, font=FONT_REGULAR):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_image(slide, path, left, top, width=None, height=None):
    if Path(path).exists():
        return slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    return None


def add_bullet_list(slide, items, left, top, width, height,
                    size=14, color=CHARCOAL, font=FONT_REGULAR):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        if item == "":
            run.text = ""
        else:
            prefix = "• " if not item.startswith(("  ", "→", "★", "·")) else ""
            run.text = f"{prefix}{item}"
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(4)
    return tb


def add_header(slide, chapter_num, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Cm(3.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    bar.shadow.inherit = False

    add_text(slide, chapter_num, Cm(1.0), Cm(0.4), Cm(6), Cm(1.0),
             size=14, bold=True, color=ACCENT, font=FONT_BOLD)
    add_text(slide, title, Cm(1.0), Cm(1.0), Cm(30), Cm(1.5),
             size=26, bold=True, color=WHITE, font=FONT_BOLD)
    if subtitle:
        add_text(slide, subtitle, Cm(1.0), Cm(2.3), Cm(30), Cm(0.7),
                 size=12, color=ICE, font=FONT_REGULAR)


def add_footer(slide, page_num):
    add_text(slide, f"{page_num} / {TOTAL}", Cm(30), Cm(18),
             Cm(3), Cm(0.8), size=10, color=MUTED, align=PP_ALIGN.RIGHT)
    add_text(slide, "CV Final · U-Net 기반 성형 견적 시각화",
             Cm(1), Cm(18), Cm(20), Cm(0.8), size=10, color=MUTED)


# ═══════════════════════════════════════
prs = make_pres()


# ━━━ 1. 표지 ━━━
s = add_blank_slide(prs)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()

add_text(s, "Computer Vision Final · 2026", Cm(2), Cm(4),
         Cm(30), Cm(0.8), size=14, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)
add_text(s, "U-Net 기반 성형 견적 시각화 시스템",
         Cm(2), Cm(6.5), Cm(30), Cm(2.5),
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_BOLD)
add_text(s, "Multi-stage U-Net Enhancement with Landmark-guided Heatmap",
         Cm(2), Cm(9.5), Cm(30), Cm(1.2),
         size=18, color=ICE, align=PP_ALIGN.CENTER)
add_text(s, "2026-06-15",
         Cm(2), Cm(13), Cm(30), Cm(1), size=14, color=ICE, align=PP_ALIGN.CENTER)
add_text(s, "kim · 팀원1 · 팀원2",
         Cm(2), Cm(14.5), Cm(30), Cm(1), size=14, color=ICE, align=PP_ALIGN.CENTER)


# ━━━ 2. 목차 ━━━
s = add_blank_slide(prs)
add_header(s, "CONTENTS", "목차", "U-Net 강화 중심 발표")

chapters = [
    ("01.", "문제 정의 + 데이터셋", "성형 견적 시각화 · CelebAMask-HQ"),
    ("02.", "시스템 + U-Net 모델 구조", "ResNet-34 Encoder · SCSE Attention"),
    ("03.", "U-Net 강화 방법론", "Combo Loss · LM Heatmap · Attention · TTA"),
    ("04.", "성능 분석 + 시각화", "Confusion Matrix · Grad-CAM · Feature Map"),
    ("05.", "디스커션 + 한계 + 향후", "어떤 강화가 가장 효과적이었나"),
]
y = Cm(4.5)
for num, title, sub in chapters:
    add_text(s, num, Cm(3), y, Cm(3), Cm(1.5), size=28, bold=True, color=ACCENT)
    add_text(s, title, Cm(6), y, Cm(20), Cm(1.0), size=20, bold=True, color=NAVY)
    add_text(s, sub, Cm(6), y + Cm(1.0), Cm(25), Cm(0.7), size=12, color=MUTED)
    y += Cm(2.4)
add_footer(s, 2)


# ━━━ 3. 문제 정의 + 데이터셋 ━━━
s = add_blank_slide(prs)
add_header(s, "01 · 문제 정의 + 데이터", "왜 U-Net 기반 시술 분석인가?")

# 왼쪽: 문제 정의
add_text(s, "문제",
         Cm(1), Cm(3.5), Cm(15), Cm(1), size=18, bold=True, color=ACCENT)
add_bullet_list(s, [
    "성형 상담 시 시술 결과 예측 불가",
    "여러 부위 동시 시술 시각화 X",
    "사용자 친화적 AI 보조 도구 부족",
    "",
    "해결: U-Net 자동 분석 + AI 시뮬레이션",
], Cm(1), Cm(4.8), Cm(15), Cm(6), size=14)

# 오른쪽: 데이터셋 통계
add_text(s, "데이터: CelebAMask-HQ (Lee 2020 CVPR)",
         Cm(17), Cm(3.5), Cm(15), Cm(1), size=18, bold=True, color=NAVY)
add_bullet_list(s, [
    "30,000장 (Train 27K / Val 3K)",
    "1024×1024 고해상도",
    "19-class 픽셀 mask annotation",
    "",
    "Task: Semantic Segmentation (6-class)",
    "  bg / nose / eye / mouth / skin / unused",
], Cm(17), Cm(4.8), Cm(15), Cm(6), size=14)

# 하단: 키 메시지
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(11.5), Cm(31), Cm(5))
box.fill.solid()
box.fill.fore_color.rgb = LIGHT_BG
box.line.color.rgb = NAVY
box.line.width = Pt(1.5)
box.shadow.inherit = False

add_text(s, "★ 본 프로젝트의 핵심",
         Cm(2), Cm(12), Cm(28), Cm(1), size=18, bold=True, color=NAVY)
add_bullet_list(s, [
    "U-Net 모델을 4단계로 강화: Baseline → LM → Attention → (TTA)",
    "수업에서 배운 CNN 원리 (계층적 추출, Skip, Attention, Loss) 모두 적용",
    "정량 + 정성 분석 (Confusion Matrix + Grad-CAM + Feature Map)",
], Cm(2), Cm(13.2), Cm(29), Cm(3.5), size=13)
add_footer(s, 3)


# ━━━ 4. 시스템 + U-Net 모델 구조 ━━━
s = add_blank_slide(prs)
add_header(s, "02 · 시스템 + 모델", "U-Net 구조 + 전체 파이프라인")

# 왼쪽: U-Net 구조 (시각적)
add_text(s, "U-Net 아키텍처 (Ronneberger 2015 MICCAI)",
         Cm(1), Cm(3.5), Cm(16), Cm(1), size=16, bold=True, color=NAVY)

# Encoder 박스
encoder_y = Cm(5)
for i, (size_label, ch) in enumerate([("256×256", "3ch"), ("128×128", "64ch"),
                                      ("64×64", "128ch"), ("32×32", "256ch"),
                                      ("16×16", "512ch")]):
    y = encoder_y + Cm(0.8) * i
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1.5), y, Cm(3.5), Cm(0.7))
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    box.shadow.inherit = False
    add_text(s, f"{size_label} · {ch}", Cm(1.5), y + Cm(0.07), Cm(3.5), Cm(0.6),
             size=10, color=WHITE, align=PP_ALIGN.CENTER, bold=True)

# Decoder 박스
for i, (size_label, ch) in enumerate([("32×32", "256ch"), ("64×64", "128ch"),
                                       ("128×128", "64ch"), ("256×256", "6 cls")]):
    y = encoder_y + Cm(0.8) * i
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(11), y, Cm(3.5), Cm(0.7))
    box.fill.solid()
    box.fill.fore_color.rgb = ACCENT
    box.line.fill.background()
    box.shadow.inherit = False
    add_text(s, f"{size_label} · {ch}", Cm(11), y + Cm(0.07), Cm(3.5), Cm(0.6),
             size=10, color=WHITE, align=PP_ALIGN.CENTER, bold=True)

# Skip connection 화살표 텍스트
add_text(s, "Skip", Cm(6), encoder_y + Cm(0.1), Cm(4), Cm(0.6),
         size=11, color=GREEN, align=PP_ALIGN.CENTER, bold=True)
add_text(s, "Connection", Cm(6), encoder_y + Cm(0.8), Cm(4), Cm(0.6),
         size=11, color=GREEN, align=PP_ALIGN.CENTER, bold=True)

# 라벨
add_text(s, "Encoder", Cm(1.5), encoder_y - Cm(0.8), Cm(4), Cm(0.7),
         size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
add_text(s, "Decoder", Cm(11), encoder_y - Cm(0.8), Cm(4), Cm(0.7),
         size=12, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# 핵심 설명
add_text(s, "구조적 특징",
         Cm(1), Cm(12.5), Cm(16), Cm(1), size=14, bold=True, color=NAVY)
add_bullet_list(s, [
    "ResNet-34 Encoder (He 2016 CVPR) — 4-stage downsample",
    "Skip Connection — multi-scale 정보 보존",
    "Decoder + SCSE Attention (Roy 2018 MICCAI)",
    "Output: (B, 6, 256, 256) — 6-class logit",
], Cm(1), Cm(13.5), Cm(16), Cm(4), size=12)

# 오른쪽: 전체 시스템 파이프라인
add_text(s, "전체 파이프라인",
         Cm(18), Cm(3.5), Cm(15), Cm(1), size=16, bold=True, color=NAVY)

stages = [
    ("1. 입력", "MediaPipe 478 landmarks"),
    ("2. 자동 입력 생성", "Bezier sketch + mask"),
    ("3. U-Net 분석 ⭐", "ResNet-34 + SCSE"),
    ("4. SD Inpaint", "Stable Diffusion 2022"),
    ("5. Refinement", "L1 + VGG19 Perceptual"),
]
y = Cm(4.7)
for step, desc in stages:
    is_unet = "U-Net" in step
    col = ACCENT if is_unet else NAVY
    add_text(s, step, Cm(18), y, Cm(8), Cm(0.7), size=12, bold=True, color=col)
    add_text(s, desc, Cm(26), y, Cm(7), Cm(0.7), size=10, color=MUTED)
    y += Cm(0.95)

add_text(s, "→ U-Net이 시술 영역 분석의 핵심",
         Cm(18), Cm(11), Cm(15), Cm(1), size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_footer(s, 4)


# ━━━ 5. U-Net 강화 방법론 ⭐ NEW ━━━
s = add_blank_slide(prs)
add_header(s, "03 · U-Net 강화 방법론", "4단계 강화 + 참고 문헌")

# 강화 방법 4개 박스 (2x2)
methods = [
    ("① Combo Loss", "Dice + CE 가중합",
     "Taghanaki 2019\nMilletari 2016 (V-Net)", NAVY),
    ("② LM-guided Heatmap", "4채널 입력 (Gaussian σ=3)",
     "Newell 2016 ECCV (Hourglass)\nKartynnik 2019 (MediaPipe)", GOLD),
    ("③ SCSE Attention", "Spatial + Channel SE",
     "Roy 2018 MICCAI\nHu 2018 CVPR (SE-Net)", NAVY),
    ("④ TTA + Early Stopping", "Hflip+Rot 앙상블 + patience=5",
     "Krizhevsky 2012 (TTA)\nPrechelt 1998 (Early Stop)", NAVY),
]

y_start = Cm(4)
x_start = Cm(1)
box_w = Cm(15.5)
box_h = Cm(5.5)

for i, (name, desc, ref, col) in enumerate(methods):
    row = i // 2
    col_idx = i % 2
    x = x_start + Cm(16.5) * col_idx
    y = y_start + Cm(6) * row

    # 박스
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = col
    box.line.width = Pt(2)
    box.shadow.inherit = False

    add_text(s, name, x + Cm(0.5), y + Cm(0.3), box_w - Cm(1), Cm(1),
             size=15, bold=True, color=col)
    add_text(s, desc, x + Cm(0.5), y + Cm(1.5), box_w - Cm(1), Cm(1),
             size=12, color=CHARCOAL)
    add_text(s, ref, x + Cm(0.5), y + Cm(3), box_w - Cm(1), Cm(2.3),
             size=11, color=MUTED)

# 강조 라벨 (② LM-guided에)
add_text(s, "⭐ 가장 효과적", Cm(28), Cm(4.3), Cm(5), Cm(0.7),
         size=10, color=ACCENT, bold=True)

add_footer(s, 5)


# ━━━ 6. U-Net 4종 비교 + 성능 ━━━
s = add_blank_slide(prs)
add_header(s, "04 · 성능 분석", "U-Net 4종 비교 (5-class mIoU)")

# 핵심 수치 박스
add_text(s,
    "Baseline 0.615  →  LM-guided 0.651 (+3.66%p ⭐)  →  Attention 0.652 (+0.03%p · plateau)",
    Cm(1), Cm(3.7), Cm(32), Cm(1.2),
    size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s,
    "※ 5-class mIoU (unused 제외, 공식 보고용) · 500장 평가",
    Cm(1), Cm(4.7), Cm(32), Cm(0.7),
    size=11, color=MUTED, align=PP_ALIGN.CENTER)

# 결과 그래프 (좌측)
img_path = SAMPLES_DIR / "model_comparison_summary.png"
add_image(s, img_path, Cm(1), Cm(5.5), width=Cm(15))

# 핵심 발견 (우측)
add_text(s, "주요 발견",
         Cm(17), Cm(5.5), Cm(15), Cm(1), size=16, bold=True, color=NAVY)
add_bullet_list(s, [
    "LM-guided 효과 (작은 영역):",
    "  eye +5.45%p ⭐ (Δ 가장 큼)",
    "  mouth +4.25%p",
    "  nose +3.62%p",
    "",
    "→ Newell 2016 ECCV 검증",
    "  (Hourglass landmark heatmap)",
    "",
    "Attention plateau (+0.03%p)",
    "→ He 2019 CVPR Bag of Tricks",
    '  "Architecture < Methodology"',
], Cm(17), Cm(6.5), Cm(16), Cm(11), size=12)
add_footer(s, 6)


# ━━━ 7. 시각화 (Grad-CAM + Feature Map 통합) ━━━
s = add_blank_slide(prs)
add_header(s, "04 · 시각화", "Grad-CAM + Feature Map · 설명 가능한 비전 AI")

# 왼쪽: Feature Map (계층적 특징)
add_text(s, "Feature Map · 계층적 특징 추출",
         Cm(1), Cm(3.5), Cm(15), Cm(1), size=14, bold=True, color=NAVY)
add_text(s, "(Zeiler 2014 ECCV)",
         Cm(1), Cm(4.2), Cm(15), Cm(0.6), size=10, color=MUTED)
img_path = SAMPLES_DIR / "feature_maps_averaged.png"
add_image(s, img_path, Cm(1), Cm(5), width=Cm(15))

add_bullet_list(s, [
    "Layer 1→4: edge → semantic 추상화",
    "256→8 해상도, 64→512 채널",
    "공간 → 의미 정보 변환 시각 증명",
], Cm(1), Cm(11), Cm(15), Cm(3), size=11)

# 오른쪽: Grad-CAM
add_text(s, "Grad-CAM · 모델이 보는 영역",
         Cm(17), Cm(3.5), Cm(15), Cm(1), size=14, bold=True, color=NAVY)
add_text(s, "(Selvaraju 2017 ICCV)",
         Cm(17), Cm(4.2), Cm(15), Cm(0.6), size=10, color=MUTED)
img_path = SAMPLES_DIR / "gradcam_presentation.png"
add_image(s, img_path, Cm(17), Cm(5), width=Cm(16))

add_bullet_list(s, [
    "'코' 클래스 → 코 영역 활성화",
    "'눈' 클래스 → 눈 영역 활성화",
    "→ 모델이 정확한 의미 학습 증명",
], Cm(17), Cm(11), Cm(15), Cm(3), size=11)

# 결론 박스
add_text(s,
    "→ 'Architecture가 의미 있게 학습됨' 시각적 증명",
    Cm(1), Cm(15.5), Cm(32), Cm(1),
    size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_footer(s, 7)


# ━━━ 8. Confusion Matrix + Failure Case ━━━
s = add_blank_slide(prs)
add_header(s, "04 · 정량 분석", "Confusion Matrix · Per-class · Failure Case")

# Confusion Matrix (좌측)
img_path = SAMPLES_DIR / "confusion_matrices_all.png"
add_image(s, img_path, Cm(1), Cm(3.8), width=Cm(18))

# 분석 결과 (우측)
add_text(s, "Per-class IoU 패턴",
         Cm(20), Cm(4), Cm(13), Cm(1), size=14, bold=True, color=NAVY)
add_bullet_list(s, [
    "쉬움: bg 0.89, skin 0.64",
    "어려움: eye 0.45 ⚠",
    "",
    "Failure Case (Worst 5):",
    "  Best 0.74 vs Worst 0.42",
    "  eye Δ=0.384 (가장 불안정)",
    "",
    "→ 측면 얼굴, 안경, 그림자",
    "→ Class imbalance 한계",
    "(Cohen 1960 · Long 2015 CVPR)",
], Cm(20), Cm(5), Cm(13), Cm(11), size=11)
add_footer(s, 8)


# ━━━ 9. SD Inpaint Before/After (샘플 1쌍만) ━━━
s = add_blank_slide(prs)
add_header(s, "U-Net + 후속 응용", "Stable Diffusion Inpaint (Rombach 2022)")

add_text(s,
    "U-Net이 분석한 시술 영역에 SD Inpaint로 변형 시뮬레이션",
    Cm(1), Cm(3.5), Cm(32), Cm(1),
    size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Before + After 1쌍 (대표 — nose_tip)
img_before = SAMPLES_DIR / "original_before.png"
img_after = SAMPLES_DIR / "sd_final_nose_tip.png"

add_text(s, "Before (원본)",
         Cm(3), Cm(5), Cm(13), Cm(1),
         size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_image(s, img_before, Cm(3), Cm(6), width=Cm(13))

add_text(s, "After (시술 시뮬레이션)",
         Cm(18), Cm(5), Cm(13), Cm(1),
         size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_image(s, img_after, Cm(18), Cm(6), width=Cm(13))

# 하단 설명
add_bullet_list(s, [
    "원 계획 SC-FEGAN (TF 1.15) → Python 3.12 환경 비호환",
    "동일 task 더 최신 모델 Stable Diffusion Inpaint(2022)로 우회",
    "→ 본 프로젝트 모듈식 설계 강점 검증 (GAN backbone 교체 가능)",
], Cm(2), Cm(14.5), Cm(30), Cm(3), size=12)
add_footer(s, 9)


# ━━━ 10. Live 데모 ━━━
s = add_blank_slide(prs)
add_header(s, "데모", "Live 시연 · Colab T4 GPU")

add_text(s,
    "▶ Live 데모",
    Cm(1), Cm(5.5), Cm(32), Cm(2),
    size=44, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s,
    "Colab 노트북 13",
    Cm(1), Cm(8.5), Cm(32), Cm(1.2),
    size=22, color=ACCENT, align=PP_ALIGN.CENTER, bold=True)
add_text(s,
    "사진 업로드 → U-Net 시술 영역 분석 → SD Inpaint 변형 → 결과",
    Cm(1), Cm(11), Cm(32), Cm(1),
    size=16, color=CHARCOAL, align=PP_ALIGN.CENTER)
add_text(s,
    "T4 GPU 기준 약 90초",
    Cm(1), Cm(12.5), Cm(32), Cm(1),
    size=13, color=MUTED, align=PP_ALIGN.CENTER, bold=True)
add_text(s,
    "github.com/keonU206/cv-final · notebooks/13_sd_inpaint.ipynb",
    Cm(1), Cm(15), Cm(32), Cm(1),
    size=10, color=MUTED, align=PP_ALIGN.CENTER)
add_footer(s, 10)


# ━━━ 11. ⭐ 디스커션 — 어떤 강화가 가장 효과적이었나 NEW ━━━
s = add_blank_slide(prs)
add_header(s, "05 · 디스커션", "어떤 U-Net 강화가 가장 효과적이었나?")

# 표 헤더
table_y = Cm(4)
headers = ["순위", "강화 방법", "효과", "참고문헌"]
widths = [Cm(2.5), Cm(11), Cm(7.5), Cm(11.5)]

# 헤더 배경
header_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Cm(1), table_y, sum(widths) + Cm(0.5), Cm(1))
header_box.fill.solid()
header_box.fill.fore_color.rgb = NAVY
header_box.line.fill.background()
header_box.shadow.inherit = False

x = Cm(1.2)
for i, h in enumerate(headers):
    add_text(s, h, x, table_y + Cm(0.2), widths[i], Cm(0.7),
             size=13, bold=True, color=WHITE)
    x += widths[i]

# 표 row
rankings = [
    ("🥇", "LM-guided Heatmap", "+3.66%p ⭐", "Newell 2016 ECCV (Hourglass)"),
    ("🥈", "Early Stopping (학습 방법)", "+2.6%p", "Prechelt 1998 (Tricks of Trade)"),
    ("🥉", "Combo Loss (Dice + CE)", "기본 학습 안정화", "Taghanaki 2019"),
    ("4", "Transfer Learning (ImageNet)", "필수 (수렴 가속)", "He 2016 / Deng 2009"),
    ("5", "SCSE Attention", "+0.03%p (plateau)", "Roy 2018 MICCAI"),
    ("6", "TTA (Test-Time)", "보합", "Krizhevsky 2012 / Shanmugam 2021"),
]

for i, (rank, method, effect, ref) in enumerate(rankings):
    y = table_y + Cm(1) + Cm(1.1) * i
    bg_color = LIGHT_BG if i % 2 == 0 else WHITE
    row_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Cm(1), y, sum(widths) + Cm(0.5), Cm(1.05))
    row_box.fill.solid()
    row_box.fill.fore_color.rgb = bg_color
    row_box.line.color.rgb = MUTED
    row_box.line.width = Pt(0.25)
    row_box.shadow.inherit = False

    x = Cm(1.2)
    cells = [rank, method, effect, ref]
    colors = [ACCENT if "🥇" in rank or "🥈" in rank or "🥉" in rank else CHARCOAL,
              NAVY, ACCENT, MUTED]
    sizes = [13, 12, 12, 10]
    bolds = [True, True, True, False]
    for j, (cell, c, sz, bd) in enumerate(zip(cells, colors, sizes, bolds)):
        add_text(s, cell, x, y + Cm(0.2), widths[j], Cm(0.7),
                 size=sz, color=c, bold=bd)
        x += widths[j]

# 결론 박스
y_conclusion = Cm(12)
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), y_conclusion, Cm(31), Cm(4.5))
box.fill.solid()
box.fill.fore_color.rgb = LIGHT_BG
box.line.color.rgb = ACCENT
box.line.width = Pt(2)
box.shadow.inherit = False

add_text(s, "★ 핵심 결론",
         Cm(2), y_conclusion + Cm(0.3), Cm(28), Cm(1),
         size=16, bold=True, color=ACCENT)
add_bullet_list(s, [
    "LM-guided Heatmap이 가장 효과적 — 특히 작은 영역(eye +5.45%p)",
    'Attention plateau 발견 → "Architecture < Training Methodology" (He 2019 CVPR Bag of Tricks)',
    "→ 단순 아키텍처 확장보다 학습 방법(Early Stopping)이 더 큰 영향",
], Cm(2), y_conclusion + Cm(1.4), Cm(29), Cm(3), size=12)
add_footer(s, 11)


# ━━━ 12. 한계 + 향후 + 결론 ━━━
s = add_blank_slide(prs)
add_header(s, "05 · 마무리", "한계 + 향후 강화 + 결론")

# 현재 한계 (좌측)
add_text(s, "현재 한계",
         Cm(1), Cm(3.5), Cm(15), Cm(1), size=16, bold=True, color=ACCENT)
add_bullet_list(s, [
    "U-Net mIoU plateau (5-class 0.65)",
    "  → Attention 이상의 효과 X",
    "",
    "eye 클래스 가장 불안정 (Δ=0.38)",
    "  → 좌우 분리 + 안경 outlier",
    "",
    "한국인 얼굴 특화 미적용",
    "  → CelebA 일반 분포만 학습",
], Cm(1), Cm(4.7), Cm(15), Cm(7), size=12)

# 향후 강화 방향 (우측)
add_text(s, "향후 강화 방향",
         Cm(17), Cm(3.5), Cm(15), Cm(1), size=16, bold=True, color=NAVY)
add_bullet_list(s, [
    "Data Augmentation 강화",
    "  CutMix (Yun 2019 ICCV)",
    "  Hard Example Mining (Lin 2017)",
    "",
    "GAN 측면 강화",
    "  ControlNet (Zhang 2023 ICCV)",
    "  InstantID (Wang 2024)",
    "",
    "Domain Adaptation",
    "  Korean Face LoRA (Hu 2022)",
], Cm(17), Cm(4.7), Cm(15), Cm(7), size=12)

# 결론 박스
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(12.5), Cm(31), Cm(5))
box.fill.solid()
box.fill.fore_color.rgb = LIGHT_BG
box.line.color.rgb = NAVY
box.line.width = Pt(2)
box.shadow.inherit = False

add_text(s, "🎯 결론",
         Cm(2), Cm(13), Cm(29), Cm(1), size=18, bold=True, color=NAVY)
add_bullet_list(s, [
    "U-Net 4단계 강화 + 정량 평가 (Confusion Matrix + Failure Case)",
    "LM-guided Heatmap이 가장 효과적 — eye/mouth 작은 영역에 +5.45%p / +4.25%p",
    'Grad-CAM + Feature Map으로 "설명 가능한 비전 AI" 만족',
    "모듈식 설계로 GAN backbone(SD Inpaint) 교체 가능 — model-agnostic",
], Cm(2), Cm(14.2), Cm(29), Cm(3), size=12)
add_footer(s, 12)


# ━━━ 13. 감사합니다 ━━━
s = add_blank_slide(prs)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()

add_text(s, "감사합니다",
         Cm(1), Cm(6.5), Cm(32), Cm(3),
         size=72, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font=FONT_BOLD)
add_text(s, "Thank You",
         Cm(1), Cm(10), Cm(32), Cm(1.5),
         size=24, color=ICE, align=PP_ALIGN.CENTER)
add_text(s, "Q & A",
         Cm(1), Cm(13), Cm(32), Cm(1.5),
         size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, "github.com/keonU206/cv-final",
         Cm(1), Cm(16), Cm(32), Cm(0.8),
         size=12, color=ICE, align=PP_ALIGN.CENTER)


# ─── 저장 ───
prs.save(str(PPTX_PATH))
print(f"✅ 발표 슬라이드 v5 생성 완료: {PPTX_PATH}")
print(f"   슬라이드 수: {len(prs.slides)}")
import os
print(f"   파일 크기: {os.path.getsize(PPTX_PATH) / 1024:.1f} KB")
print()
print("📋 v5 구성 (13장, U-Net 강화 중심):")
print("  1.  표지")
print("  2.  목차 (5 챕터)")
print("  3.  문제 정의 + 데이터셋 (통합)")
print("  4.  시스템 + U-Net 모델 구조 ⭐ (모델 구조 시각화 추가)")
print("  5.  U-Net 강화 방법론 ⭐ NEW (4가지 방법 + 참고문헌)")
print("  6.  U-Net 4종 비교 + 성능")
print("  7.  시각화 (Grad-CAM + Feature Map 통합)")
print("  8.  Confusion Matrix + Failure Case (통합)")
print("  9.  SD Inpaint Before/After (1쌍만)")
print("  10. Live 데모")
print("  11. 디스커션 ⭐ NEW — 어떤 강화가 가장 효과적이었나")
print("  12. 한계 + 향후 + 결론")
print("  13. 감사합니다")
