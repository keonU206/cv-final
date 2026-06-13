"""시계열_4조.pptx 분석 (큰 파일)."""
import os
from pathlib import Path

os.environ["PYTHONIOENCODING"] = "utf-8"

# Python pathlib 사용 — 한글 경로 안전
folder = Path(r"C:\Users\User\Documents\카카오톡 받은 파일")
for f in folder.iterdir():
    if "시계열_4조" in f.name and f.suffix == ".pptx":
        print(f"파일 발견: {f.name} ({f.stat().st_size / (1024*1024):.1f} MB)")
        try:
            from pptx import Presentation
            prs = Presentation(str(f))
            w_in = prs.slide_width / 914400
            h_in = prs.slide_height / 914400
            print(f"크기: {w_in:.1f} × {h_in:.1f} (비율 {w_in/h_in:.2f})")
            print(f"슬라이드 수: {len(prs.slides)}")

            for i, slide in enumerate(prs.slides, 1):
                title = ""
                image_count = 0
                video_count = 0
                body_texts = []

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                text = run.text.strip()
                                if text:
                                    if not title and 2 < len(text) < 80:
                                        title = text
                                    elif len(text) > 5:
                                        body_texts.append(text[:50])
                    if shape.shape_type == 13:  # PICTURE
                        image_count += 1
                    # 비디오는 shape_type == 17 (MEDIA) 또는 element 안에 video
                    if hasattr(shape, "element"):
                        if "video" in str(shape.element.xml).lower():
                            video_count += 1

                marks = []
                if image_count: marks.append(f"🖼{image_count}")
                if video_count: marks.append(f"🎬{video_count}")
                mark_str = " ".join(marks) if marks else ""
                title_s = (title[:50] + "...") if len(title) > 50 else title
                print(f"  {i:2}. {title_s:<55} {mark_str}")
                for t in body_texts[:2]:
                    if t != title:
                        print(f"      → {t}")
        except Exception as e:
            print(f"❌ 에러: {e}")
        break
