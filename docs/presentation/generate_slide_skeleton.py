"""발표 슬라이드 틀 PPTX (skeleton) 생성 — 11장.

- 텍스트는 모두 포함
- 이미지 자리는 회색 placeholder 박스 + 가이드 텍스트
- 가이드 형식: "📁 노트북 N, 셀 #M 결과 이미지 첨부 부탁 (파일명)"

출력: 바탕화면/발표자료_CV_Final/발표슬라이드_틀.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt

# ─── 색상 ───
NAVY = RGBColor(0x1E, 0x27, 0x61)
ACCENT = RGBColor(0xF9, 0x61, 0x67)
CHARCOAL = RGBColor(0x2C, 0x3E, 0x50)
MUTED = RGBColor(0x7B, 0x8F, 0xA1)
LIGHT_GRAY = RGBColor(0xE2, 0xE8, 0xF0)
BG = RGBColor(0xFA, 0xFB, 0xFD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xF5, 0xA6, 0x23)
GREEN = RGBColor(0x48, 0xBB, 0x78)
PLACEHOLDER_FILL = RGBColor(0xF0, 0xF4, 0xF8)
PLACEHOLDER_BORDER = RGBColor(0xCB, 0xD5, 0xE0)

FONT = "맑은 고딕"

# ─── 16:9 ───
prs = Presentation()
prs.slide_width = Cm(33.867)
prs.slide_height = Cm(19.05)
BLANK = prs.slide_layouts[6]


def add_text(slide, x, y, w, h, text, *, font_size=14, bold=False,
             color=CHARCOAL, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT):
    """텍스트 박스 추가."""
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill, line=None, line_width=0):
    """배경 사각형."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Cm(x), Cm(y), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    shape.shadow.inherit = False
    return shape


def add_round_rect(slide, x, y, w, h, fill, line=None, line_width=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Cm(x), Cm(y), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    return shape


def add_image_placeholder(slide, x, y, w, h, notebook, cell, filename,
                          description=""):
    """이미지 자리 placeholder 박스 + 가이드 텍스트."""
    # 배경 박스 (점선처럼 보이게 두꺼운 테두리)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Cm(x), Cm(y), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PLACEHOLDER_FILL
    shape.line.color.rgb = PLACEHOLDER_BORDER
    shape.line.width = Pt(2)

    # 가이드 텍스트
    guide = f"📁 이미지 자리\n\n"
    guide += f"노트북 {notebook} · 셀 #{cell}\n"
    guide += f"결과 이미지 첨부 부탁\n\n"
    guide += f"📄 {filename}"
    if description:
        guide += f"\n\n({description})"

    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Cm(0.3)
    tf.margin_right = Cm(0.3)

    lines = guide.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(11)
        if "📁" in line or "📄" in line:
            run.font.bold = True
            run.font.color.rgb = ACCENT
        else:
            run.font.color.rgb = MUTED


def slide_header(slide, chapter, title, subtitle=""):
    """공통 슬라이드 헤더 (상단 챕터 + 제목)."""
    # 좌측 컬러 바
    add_rect(slide, 0, 0, 0.5, 19.05, NAVY)

    # 챕터 번호
    add_text(slide, 1.2, 0.8, 5, 0.8, chapter,
             font_size=11, bold=True, color=ACCENT)
    # 제목
    add_text(slide, 1.2, 1.4, 30, 1.2, title,
             font_size=24, bold=True, color=NAVY)
    # 부제목
    if subtitle:
        add_text(slide, 1.2, 2.6, 30, 0.6, subtitle,
                 font_size=12, color=MUTED)

    # 하단 페이지 표시 (슬라이드 번호는 따로 추가)


def slide_footer(slide, page_num, total=11):
    """하단 페이지 번호."""
    add_text(slide, 30, 18, 3, 0.5, f"{page_num} / {total}",
             font_size=9, color=MUTED, align=PP_ALIGN.RIGHT)
    add_text(slide, 1.2, 18, 20, 0.5,
             "U-Net 기반 성형 견적 시각화 시스템 · 2026-06-15",
             font_size=9, color=MUTED)


# ═══════════════════════════════════════════════════════════════
# Slide 1 · 표지
# ═══════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
# 배경
add_rect(s1, 0, 0, 33.867, 19.05, NAVY)
# 상단 액센트 바
add_rect(s1, 0, 8, 33.867, 0.15, ACCENT)

# 메인 제목
add_text(s1, 2, 5.5, 30, 2, "U-Net 기반 성형 견적",
         font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s1, 2, 7.3, 30, 2, "시각화 시스템",
         font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 부제목
add_text(s1, 2, 10, 30, 1,
         "Multi-stage U-Net Enhancement with Landmark-guided Heatmap",
         font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 발표 정보 박스
add_round_rect(s1, 8, 13.5, 17.867, 3, RGBColor(0x2A, 0x35, 0x7A))
add_text(s1, 8, 13.8, 17.867, 0.7,
         "Computer Vision Final Project · 2026-06-15",
         font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s1, 8, 14.6, 17.867, 0.7,
         "[팀명] · 발표자: A · B · C",
         font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s1, 8, 15.4, 17.867, 0.7,
         "지도교수: 박한샘",
         font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# Slide 2 · 목차
# ═══════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, 33.867, 19.05, BG)
slide_header(s2, "TABLE OF CONTENTS", "목차",
             "10분 발표 · 5개 챕터로 진행")

chapters = [
    ("01", "문제 정의 + 데이터셋", "Slide 3~4"),
    ("02", "시스템 + U-Net 모델 구조 + 선택 근거", "Slide 5~6"),
    ("03", "U-Net 강화 방법론 + 결과", "Slide 7"),
    ("04", "시각화 + 해석 (Grad-CAM + Feature Map)", "Slide 8, 10"),
    ("05", "Live 시연 + 디스커션 + 마무리", "Slide 9~11"),
]

y_start = 4.5
for i, (ch_num, ch_title, ch_pages) in enumerate(chapters):
    y = y_start + i * 2.3
    # 챕터 번호
    add_round_rect(s2, 2, y, 2, 1.8, NAVY)
    add_text(s2, 2, y, 2, 1.8, ch_num,
             font_size=24, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 제목
    add_text(s2, 5, y + 0.2, 22, 0.9, ch_title,
             font_size=18, bold=True, color=NAVY)
    add_text(s2, 5, y + 1.1, 22, 0.7, ch_pages,
             font_size=11, color=MUTED)

slide_footer(s2, 2)


# ═══════════════════════════════════════════════════════════════
# Slide 3 · 문제 정의
# ═══════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, 33.867, 19.05, BG)
slide_header(s3, "01 · PROBLEM DEFINITION", "문제 정의",
             "Task: Semantic Segmentation + Image Generation")

# 좌측 — 사용자 어려움
add_round_rect(s3, 1.5, 4.5, 14.5, 11.5, WHITE,
               line=LIGHT_GRAY, line_width=1)
add_text(s3, 2, 5, 13.5, 0.8, "🚨 현재 사용자 어려움",
         font_size=16, bold=True, color=ACCENT)

problems = [
    "1. 상담 전 시술 결과 예측 불가",
    "2. 여러 부위 동시 시술 시각화 X",
    "3. 병원별 견적 비교 어려움",
    "4. AI 보조 도구 부족",
]
for i, p in enumerate(problems):
    add_text(s3, 2.3, 6.5 + i * 2, 13, 1.5, p,
             font_size=14, color=CHARCOAL)

# 우측 — 해결
add_round_rect(s3, 17.5, 4.5, 14.5, 11.5, WHITE,
               line=LIGHT_GRAY, line_width=1)
add_text(s3, 18, 5, 13.5, 0.8, "✅ 본 프로젝트 해결",
         font_size=16, bold=True, color=GREEN)

solutions = [
    "1. 사진 1장 → AI 자동 분석",
    "2. 시술 부위 자동 검출 (U-Net)",
    "3. 변형 시뮬레이션 (SD Inpaint)",
    "4. PDF 견적서 자동 생성",
]
for i, p in enumerate(solutions):
    add_text(s3, 18.3, 6.5 + i * 2, 13, 1.5, p,
             font_size=14, color=CHARCOAL)

# 핵심 메시지
add_round_rect(s3, 1.5, 16.5, 30.5, 1.2, NAVY)
add_text(s3, 1.5, 16.5, 30.5, 1.2,
         "💡 사진 1장 → AI 시뮬레이션 + 견적서 (자동화)",
         font_size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

slide_footer(s3, 3)


# ═══════════════════════════════════════════════════════════════
# Slide 4 · 사용한 데이터
# ═══════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, 33.867, 19.05, BG)
slide_header(s4, "01 · DATASET", "사용한 데이터",
             "CelebAMask-HQ (Lee 2020 CVPR · MaskGAN)")

# 통계 박스 3개
stats = [
    ("30,000장", "이미지 수", "Train 27K / Val 3K", NAVY),
    ("1024×1024", "해상도", "고해상도 face", ACCENT),
    ("19-class", "Mask Annotation", "Pixel-level", GOLD),
]
for i, (val, label, desc, color) in enumerate(stats):
    x = 1.5 + i * 10.5
    add_round_rect(s4, x, 4.5, 9.5, 5, color)
    add_text(s4, x, 5, 9.5, 1.5, val,
             font_size=28, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s4, x, 6.8, 9.5, 0.8, label,
             font_size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s4, x, 7.8, 9.5, 0.8, desc,
             font_size=11, color=LIGHT_GRAY,
             align=PP_ALIGN.CENTER)

# 활용 방식
add_round_rect(s4, 1.5, 10.5, 30.5, 6.5, WHITE,
               line=LIGHT_GRAY, line_width=1)
add_text(s4, 2, 10.8, 30, 0.8, "📊 본 프로젝트 활용",
         font_size=14, bold=True, color=NAVY)

usage = [
    "• 6-class로 재라벨: background, nose, eye, mouth, skin, unused",
    "• 5-class 평가 (unused 제외, 공식 보고용) + 6-class 평균 (학술 정직성)",
    "• MediaPipe Face Mesh (Kartynnik 2019)로 478 landmark 추가 추출",
    "• Train 27K + Val 3K → 본 실험은 Val 500장 sampling (10x 확대, 통계 안정성)",
]
for i, u in enumerate(usage):
    add_text(s4, 2.3, 11.7 + i * 1.2, 30, 1.0, u,
             font_size=12, color=CHARCOAL)

slide_footer(s4, 4)


# ═══════════════════════════════════════════════════════════════
# Slide 5 · 파이프라인 + U-Net 모델 구조
# ═══════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, 33.867, 19.05, BG)
slide_header(s5, "02 · SYSTEM PIPELINE", "전체 파이프라인 + U-Net 모델 구조",
             "7-stage modular system · Encoder-Decoder + Skip + SCSE")

# 좌측 — 7단계 파이프라인
add_round_rect(s5, 1.5, 4.5, 15.5, 12.5, WHITE,
               line=LIGHT_GRAY, line_width=1)
add_text(s5, 2, 4.8, 14.5, 0.8, "🔄 7단계 모듈식 파이프라인",
         font_size=14, bold=True, color=ACCENT)

steps = [
    ("1", "이미지 입력 + 전처리", "256×256 OpenCV"),
    ("2", "MediaPipe Face Mesh", "478 landmark"),
    ("3", "자동 입력 생성", "Mask + Sketch + Color"),
    ("4", "U-Net 분석 ⭐", "6-class segmentation"),
    ("5", "SD Inpaint", "Rombach 2022"),
    ("6", "Refinement Network", "Johnson 2016"),
    ("7", "결과 출력", "PDF + 시각화"),
]
for i, (num, title, desc) in enumerate(steps):
    y = 6 + i * 1.5
    # 원형 번호
    add_round_rect(s5, 2, y, 1, 1, NAVY if num != "4" else ACCENT)
    add_text(s5, 2, y, 1, 1, num,
             font_size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s5, 3.5, y + 0.05, 8, 0.5, title,
             font_size=12, bold=True, color=NAVY)
    add_text(s5, 3.5, y + 0.55, 12, 0.5, desc,
             font_size=10, color=MUTED)

# 우측 — U-Net 구조 도식
add_round_rect(s5, 17.5, 4.5, 14.867, 12.5, WHITE,
               line=LIGHT_GRAY, line_width=1)
add_text(s5, 18, 4.8, 14, 0.8, "🧠 U-Net 모델 구조",
         font_size=14, bold=True, color=ACCENT)
add_text(s5, 18, 5.5, 14, 0.6,
         "Ronneberger 2015 MICCAI · He 2016 · Roy 2018",
         font_size=10, color=MUTED)

# Encoder
add_round_rect(s5, 18, 6.5, 6, 1, NAVY)
add_text(s5, 18, 6.5, 6, 1, "Encoder (ResNet-34)",
         font_size=11, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 화살표 ↓
add_text(s5, 18, 7.6, 6, 0.5, "↓ downsample 4-stage",
         font_size=9, color=MUTED, align=PP_ALIGN.CENTER)

# Bottleneck
add_round_rect(s5, 19, 8.3, 4, 0.8, GOLD)
add_text(s5, 19, 8.3, 4, 0.8, "Bottleneck (512ch)",
         font_size=10, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Skip 화살표
add_text(s5, 25, 7.4, 7, 0.6, "Skip Connection →",
         font_size=10, color=ACCENT, bold=True)

# Decoder
add_round_rect(s5, 18, 9.5, 6, 1, ACCENT)
add_text(s5, 18, 9.5, 6, 1, "Decoder + SCSE",
         font_size=11, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 화살표 ↓
add_text(s5, 18, 10.6, 6, 0.5, "↓ upsample 4-stage",
         font_size=9, color=MUTED, align=PP_ALIGN.CENTER)

# Output
add_round_rect(s5, 18, 11.3, 6, 1, NAVY)
add_text(s5, 18, 11.3, 6, 1, "Output (6, 256, 256)",
         font_size=11, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 설명
add_text(s5, 18, 12.8, 14, 0.6,
         "📐 입력: (3, 256, 256) → 출력: 6-class logit",
         font_size=11, color=CHARCOAL)
add_text(s5, 18, 13.5, 14, 0.6,
         "⚡ Loss: Combo (Dice + CE) + Cosine Annealing",
         font_size=11, color=CHARCOAL)
add_text(s5, 18, 14.2, 14, 0.6,
         "🔑 Skip Connection: multi-scale 정보 보존",
         font_size=11, color=CHARCOAL)
add_text(s5, 18, 14.9, 14, 0.6,
         "🎯 SCSE Attention: 공간 + 채널 가중치",
         font_size=11, color=CHARCOAL)

# 핵심 메시지
add_text(s5, 18, 16, 14, 0.7,
         "→ U-Net이 시술 영역 분석의 핵심",
         font_size=12, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER)

slide_footer(s5, 5)


# ═══════════════════════════════════════════════════════════════
# Slide 6 · 왜 U-Net? (모델 선택 근거) — 🆕
# ═══════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, 33.867, 19.05, BG)
slide_header(s6, "02 · MODEL SELECTION", "왜 U-Net인가? (모델 선택 근거)",
             "Semantic Segmentation 6종 비교 · 학습 가능 + 강화 가능")

# 좌측 — 비교표
add_round_rect(s6, 1.5, 4.5, 17, 12.5, WHITE,
               line=LIGHT_GRAY, line_width=1)
add_text(s6, 2, 4.8, 16, 0.8, "📊 Semantic Segmentation 6종 비교",
         font_size=14, bold=True, color=ACCENT)

# 표 헤더
add_rect(s6, 2, 5.8, 16, 0.8, NAVY)
header_cols = [("모델", 2, 4), ("출처 (년도)", 6, 5),
               ("핵심", 11, 4.5), ("적합", 15.5, 2.5)]
for label, x, w in header_cols:
    add_text(s6, x, 5.8, w, 0.8, label,
             font_size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 표 데이터
models = [
    ("FCN", "Long 2015", "원조 segmentation", "❌"),
    ("SegNet", "Badrinarayanan 2017", "Skip 없음", "❌"),
    ("DeepLab v3+", "Chen 2018 ECCV", "Atrous + ASPP", "△"),
    ("U-Net ⭐", "Ronneberger 2015", "Skip + Symmetric", "✅"),
    ("HRNet", "Wang 2020", "Multi-resolution", "△"),
    ("SAM", "Kirillov 2023", "Foundation (frozen)", "❌"),
]
for i, (m, src, key, fit) in enumerate(models):
    y = 6.6 + i * 1.1
    bg_color = BG if i % 2 == 0 else WHITE
    is_unet = "U-Net" in m
    if is_unet:
        bg_color = RGBColor(0xFF, 0xF4, 0xE6)  # 강조
    add_rect(s6, 2, y, 16, 1.1, bg_color)

    text_color = ACCENT if is_unet else CHARCOAL
    add_text(s6, 2, y, 4, 1.1, m,
             font_size=11, bold=is_unet, color=text_color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s6, 6, y, 5, 1.1, src,
             font_size=10, color=text_color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s6, 11, y, 4.5, 1.1, key,
             font_size=10, color=text_color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s6, 15.5, y, 2.5, 1.1, fit,
             font_size=14, bold=True, color=text_color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 우측 — 4대 근거
add_round_rect(s6, 19.5, 4.5, 12.867, 12.5, WHITE,
               line=LIGHT_GRAY, line_width=1)
add_text(s6, 20, 4.8, 12, 0.8, "🎯 U-Net 선택 4대 근거",
         font_size=14, bold=True, color=ACCENT)

reasons = [
    ("①", "Skip Connection",
     "eye, nose 같은\n작은 영역 분할 필수"),
    ("②", "Medical/Face Parsing 표준",
     "CelebAMask-HQ\nbaseline도 U-Net 계열"),
    ("③", "소량~중간 데이터에 강함",
     "30K로 충분 학습\n(의료영상 출발)"),
    ("④", "모듈식 확장 가능 ⭐",
     "본 프로젝트 4단계\n강화의 핵심"),
]
for i, (num, title, desc) in enumerate(reasons):
    y = 6 + i * 2.6
    add_round_rect(s6, 20, y, 1.5, 1.5, NAVY)
    add_text(s6, 20, y, 1.5, 1.5, num,
             font_size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s6, 21.8, y, 10, 0.7, title,
             font_size=12, bold=True, color=NAVY)
    add_text(s6, 21.8, y + 0.7, 10, 1.3, desc,
             font_size=10, color=CHARCOAL)

# 하단 결론
add_round_rect(s6, 1.5, 17.5, 30.867, 1, ACCENT)
add_text(s6, 1.5, 17.5, 30.867, 1,
         "💡 \"학습 가능 + 강화 가능\"한 모델 = U-Net 최적 선택",
         font_size=13, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

slide_footer(s6, 6)


# ═══════════════════════════════════════════════════════════════
# Slide 7 · U-Net 강화 방법론 + 결과 ⭐ (이미지 2개)
# ═══════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
add_rect(s7, 0, 0, 33.867, 19.05, BG)
slide_header(s7, "03 · U-NET ENHANCEMENT",
             "U-Net 강화 방법론 + 결과",
             "4가지 강화 · 5-class mIoU 0.65 (공식 보고)")

# 상단 좌측 — 4가지 강화 방법
add_round_rect(s7, 1.5, 4.5, 15, 6, WHITE,
               line=LIGHT_GRAY, line_width=1)
add_text(s7, 2, 4.8, 14, 0.7, "🛠 4가지 강화 방법",
         font_size=13, bold=True, color=ACCENT)

methods = [
    ("①", "Combo Loss (Dice + CE)", "Taghanaki 2019"),
    ("②", "LM-guided Heatmap ⭐", "Newell 2016 ECCV"),
    ("③", "SCSE Attention", "Roy 2018 MICCAI"),
    ("④", "TTA + Early Stopping", "Krizhevsky 2012"),
]
for i, (num, name, ref) in enumerate(methods):
    y = 5.7 + i * 1.1
    add_text(s7, 2.3, y, 1, 0.8, num,
             font_size=14, bold=True, color=ACCENT)
    add_text(s7, 3.3, y, 8, 0.5, name,
             font_size=11, bold=True, color=NAVY)
    add_text(s7, 3.3, y + 0.5, 8, 0.5, ref,
             font_size=9, color=MUTED)

# 상단 우측 — 핵심 수치
add_round_rect(s7, 17, 4.5, 15.367, 6, NAVY)
add_text(s7, 17, 4.8, 15, 0.7, "⭐ 핵심 수치",
         font_size=13, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)

# 수치
add_text(s7, 17, 5.8, 15.367, 1, "Baseline → LM → Attention",
         font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s7, 17, 6.8, 15.367, 1.5,
         "0.615  →  0.651  →  0.652",
         font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s7, 17, 8.3, 15.367, 0.8,
         "(+3.66%p ⭐ LM-guided가 가장 효과적)",
         font_size=11, color=GOLD, align=PP_ALIGN.CENTER)
add_text(s7, 17, 9, 15.367, 0.6,
         "※ 5-class mIoU (공식 보고용, unused 제외)",
         font_size=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s7, 17, 9.6, 15.367, 0.6,
         "※ 6-class 평균: 0.54 (학술 정직성)",
         font_size=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 하단 — 이미지 2개 placeholder
add_image_placeholder(s7, 1.5, 11, 15, 6,
                      notebook="15",
                      cell="#8",
                      filename="model_comparison_summary.png",
                      description="3 모델 mIoU 비교 막대그래프")

add_image_placeholder(s7, 17, 11, 15.367, 6,
                      notebook="15",
                      cell="#5",
                      filename="confusion_matrices_all.png",
                      description="3 모델 5×5 Confusion Matrix Heatmap")

slide_footer(s7, 7)


# ═══════════════════════════════════════════════════════════════
# Slide 8 · 입력 → 결과 + Grad-CAM (이미지 3개)
# ═══════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
add_rect(s8, 0, 0, 33.867, 19.05, BG)
slide_header(s8, "04 · VISUALIZATION + EXPLAINABILITY",
             "입력 → 결과 + Grad-CAM",
             "설명 가능한 비전 AI · Selvaraju 2017 ICCV")

# 좌측 — Before/After
add_text(s8, 1.5, 4.3, 15, 0.6, "📸 Before / After (코끝 시술)",
         font_size=13, bold=True, color=NAVY)

add_image_placeholder(s8, 1.5, 5, 7, 7,
                      notebook="13",
                      cell="#5",
                      filename="original_before.png",
                      description="원본 사진")

add_image_placeholder(s8, 9, 5, 7.5, 7,
                      notebook="13",
                      cell="#7",
                      filename="sd_final_nose_tip.png",
                      description="SD Inpaint 시술 후")

# 좌측 하단 메시지
add_round_rect(s8, 1.5, 12.3, 15, 1.7, NAVY)
add_text(s8, 1.5, 12.3, 15, 1.7,
         "→ Stable Diffusion Inpaint\n   (Rombach 2022 CVPR)",
         font_size=11, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

# 우측 — Grad-CAM
add_text(s8, 17, 4.3, 15, 0.6, "🔥 Grad-CAM (모델이 어디를 보는가?)",
         font_size=13, bold=True, color=ACCENT)

add_image_placeholder(s8, 17, 5, 15.367, 7,
                      notebook="14",
                      cell="#5",
                      filename="gradcam_presentation.png",
                      description="4 클래스 (nose, eye, mouth, skin) Grad-CAM")

# 우측 하단 메시지
add_round_rect(s8, 17, 12.3, 15.367, 1.7, ACCENT)
add_text(s8, 17, 12.3, 15.367, 1.7,
         "→ '코' 클래스 → 코 영역, '눈' → 눈 영역만 활성화\n"
         "   모델이 의미 있는 semantic 학습 ✓",
         font_size=11, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 하단 — 학술 인용
add_text(s8, 1.5, 14.5, 30.867, 0.7,
         "Selvaraju 2017 ICCV — Grad-CAM | "
         "Johnson 2016 ECCV — Perceptual Loss | "
         "Rombach 2022 CVPR — Stable Diffusion",
         font_size=10, color=MUTED, align=PP_ALIGN.CENTER)

slide_footer(s8, 8)


# ═══════════════════════════════════════════════════════════════
# Slide 9 · Live 시연
# ═══════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
add_rect(s9, 0, 0, 33.867, 19.05, BG)
slide_header(s9, "05 · LIVE DEMO", "Live 시연 (Colab T4 GPU)",
             "노트북 13 · 약 90초 inference")

# 중앙 — 큰 버튼 형태
add_round_rect(s9, 6.5, 5, 21, 9, NAVY)

add_text(s9, 6.5, 5.5, 21, 1.5, "▶  Live Demo",
         font_size=48, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s9, 6.5, 8, 21, 1, "Colab Notebook 13 - SD Inpaint",
         font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 흐름
add_text(s9, 6.5, 10, 21, 0.8,
         "📷 사진 업로드  →  🧠 U-Net 분석  →  🎨 SD Inpaint  →  📊 결과",
         font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

# 시간
add_round_rect(s9, 12, 12, 10, 1.5, GOLD)
add_text(s9, 12, 12, 10, 1.5, "⏱ 약 90초 (T4 GPU)",
         font_size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 하단 안내
add_text(s9, 1.5, 14.8, 30.867, 0.7,
         "🔗 github.com/keonU206/cv-final/blob/main/notebooks/13_sd_inpaint.ipynb",
         font_size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(s9, 1.5, 15.8, 30.867, 0.7,
         "💡 발표 1시간 전 셀 1, 2 사전 실행 (SD 모델 4.27GB 로드)",
         font_size=10, color=MUTED, align=PP_ALIGN.CENTER)

add_text(s9, 1.5, 16.5, 30.867, 0.7,
         "📹 백업: 데모 영상 mp4 준비",
         font_size=10, color=MUTED, align=PP_ALIGN.CENTER)

slide_footer(s9, 9)


# ═══════════════════════════════════════════════════════════════
# Slide 10 · 디스커션 + Feature Map (이미지 1개)
# ═══════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
add_rect(s10, 0, 0, 33.867, 19.05, BG)
slide_header(s10, "05 · DISCUSSION",
             "디스커션 + Feature Map",
             "어떤 U-Net 강화가 가장 효과적이었나?")

# 좌측 — 순위 표
add_round_rect(s10, 1.5, 4.5, 16, 10.5, WHITE,
               line=LIGHT_GRAY, line_width=1)
add_text(s10, 2, 4.8, 15, 0.8, "🏆 강화 방법 효과 순위",
         font_size=13, bold=True, color=ACCENT)

rankings = [
    ("🥇", "LM-guided Heatmap", "+3.66%p ⭐", "Newell 2016", GOLD),
    ("🥈", "Early Stopping", "+2.6%p", "Prechelt 1998", MUTED),
    ("🥉", "Combo Loss", "안정화", "Taghanaki 2019", RGBColor(0xCD, 0x7F, 0x32)),
    ("4", "Transfer Learning", "필수", "He 2016", CHARCOAL),
    ("5", "SCSE Attention", "+0.03%p", "Roy 2018", CHARCOAL),
    ("6", "TTA", "보합", "Krizhevsky 2012", CHARCOAL),
]
for i, (rank, name, effect, ref, color) in enumerate(rankings):
    y = 6 + i * 1.4
    add_text(s10, 2.3, y, 1.5, 0.9, rank,
             font_size=16, bold=True, color=color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s10, 4, y, 6, 0.5, name,
             font_size=11, bold=True, color=NAVY)
    add_text(s10, 4, y + 0.5, 6, 0.4, ref,
             font_size=9, color=MUTED)
    add_text(s10, 10, y, 6, 0.9, effect,
             font_size=12, bold=True, color=color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 우측 — Feature Map + 결론
add_text(s10, 18, 4.3, 14.367, 0.6,
         "🧠 Feature Map (CNN 계층적 추상화)",
         font_size=13, bold=True, color=ACCENT)

add_image_placeholder(s10, 18, 5, 14.367, 6,
                      notebook="16",
                      cell="#5",
                      filename="feature_maps_averaged.png",
                      description="5컬럼 Layer 1~4 채널 평균")

# 결론 박스
add_round_rect(s10, 18, 11.5, 14.367, 5.5, NAVY)
add_text(s10, 18.3, 11.8, 13.5, 0.7, "💡 핵심 결론 3가지",
         font_size=13, bold=True, color=WHITE)

conclusions = [
    "1. LM-guided가 가장 효과적 (eye +5.45%p)",
    "2. Attention plateau (+0.03%p)",
    "   → \"Architecture < Methodology\"",
    "       (He 2019 CVPR Bag of Tricks)",
    "3. 학습 방법론 > 모델 확장",
]
for i, c in enumerate(conclusions):
    add_text(s10, 18.3, 12.7 + i * 0.8, 14, 0.6, c,
             font_size=11, color=LIGHT_GRAY)

slide_footer(s10, 10)


# ═══════════════════════════════════════════════════════════════
# Slide 11 · 마무리
# ═══════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK)
add_rect(s11, 0, 0, 33.867, 19.05, BG)
slide_header(s11, "05 · CONCLUSION",
             "마무리 · 한계 · 향후 강화",
             "감사합니다 · Q&A")

# 좌측 — 한계
add_round_rect(s11, 1.5, 4.5, 14.5, 10, WHITE,
               line=LIGHT_GRAY, line_width=1)
add_text(s11, 2, 4.8, 13.5, 0.8, "⚠ 현재 한계",
         font_size=14, bold=True, color=ACCENT)

limits = [
    "• U-Net mIoU plateau (5-class 0.65)",
    "• eye 클래스 가장 불안정 (Δ=0.384)",
    "• 한국인 얼굴 특화 미적용",
    "• SC-FEGAN TF 1.15 환경 호환 X",
]
for i, t in enumerate(limits):
    add_text(s11, 2.3, 5.8 + i * 1.5, 13, 1, t,
             font_size=12, color=CHARCOAL)

# 우측 — 향후 강화
add_round_rect(s11, 17.5, 4.5, 14.867, 10, NAVY)
add_text(s11, 18, 4.8, 14, 0.8, "🚀 향후 강화 방향",
         font_size=14, bold=True, color=WHITE)

futures = [
    ("Data", "CutMix (Yun 2019)\nHard Example Mining (Lin 2017)"),
    ("GAN", "ControlNet (Zhang 2023)\nInstantID (Wang 2024)"),
    ("Domain", "Korean Face LoRA (Hu 2022)"),
]
for i, (cat, desc) in enumerate(futures):
    y = 5.8 + i * 2.5
    add_text(s11, 18.3, y, 4, 0.7, cat,
             font_size=12, bold=True, color=GOLD)
    add_text(s11, 18.3, y + 0.7, 14, 1.5, desc,
             font_size=11, color=LIGHT_GRAY)

# 하단 — 감사 인사
add_round_rect(s11, 1.5, 15, 30.867, 3, ACCENT)
add_text(s11, 1.5, 15, 30.867, 1.5, "🙏 감사합니다",
         font_size=32, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s11, 1.5, 16.5, 30.867, 1.5, "질문 받겠습니다 · Q&A",
         font_size=14, color=WHITE,
         align=PP_ALIGN.CENTER)

slide_footer(s11, 11)


# ─── 저장 ───
OUTPUT_DIR = Path("C:/Users/User/Desktop/발표자료_CV_Final")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
PPTX_PATH = OUTPUT_DIR / "발표슬라이드_틀.pptx"

prs.save(str(PPTX_PATH))
print(f"✅ PPTX 생성 완료: {PPTX_PATH}")
import os
print(f"   파일 크기: {os.path.getsize(PPTX_PATH) / 1024:.1f} KB")
print(f"   슬라이드 수: {len(prs.slides)}")
