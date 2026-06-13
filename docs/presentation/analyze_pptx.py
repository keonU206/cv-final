"""3개 발표자료 분석 — 박한샘 교수님 스타일 파악용."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

PPTX_FILES = [
    "C:/Users/User/Documents/카카오톡 받은 파일/시계열_4조.pptx",
    "C:/Users/User/Documents/카카오톡 받은 파일/시계열데이터_프로젝트_최종본.pptx",
    "C:/Users/User/Documents/카카오톡 받은 파일/한국어 AI 뉴스 도메인 Temporal-aware Hybrid Retrieval 기반 RAG 성능 최적화 연구.pptx",
]


def analyze(path):
    print(f"\n{'='*70}")
    print(f"📊 {Path(path).name}")
    print(f"{'='*70}")
    prs = Presentation(path)

    # 슬라이드 크기 / 비율
    w_in = prs.slide_width / 914400  # EMU → inches
    h_in = prs.slide_height / 914400
    ratio = w_in / h_in
    print(f"슬라이드 크기: {w_in:.1f} × {h_in:.1f} inches (비율 {ratio:.2f})")
    print(f"슬라이드 수: {len(prs.slides)}")
    print(f"파일 크기: {Path(path).stat().st_size / (1024*1024):.1f} MB")

    # 슬라이드별 정보
    print(f"\n--- 슬라이드별 ---")
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        text_count = 0
        image_count = 0
        chart_count = 0
        table_count = 0
        body_texts = []

        for shape in slide.shapes:
            # 텍스트
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        text = run.text.strip()
                        if text:
                            text_count += 1
                            if not title and len(text) > 2 and len(text) < 80:
                                title = text
                            elif len(text) > 5:
                                body_texts.append(text[:60])
            # 이미지
            if shape.shape_type == 13:  # PICTURE
                image_count += 1
            # 차트
            if shape.has_chart:
                chart_count += 1
            # 표
            if shape.has_table:
                table_count += 1

        title_short = title[:50] + "..." if len(title) > 50 else title
        marks = []
        if image_count: marks.append(f"🖼{image_count}")
        if chart_count: marks.append(f"📊{chart_count}")
        if table_count: marks.append(f"📋{table_count}")
        mark_str = " ".join(marks) if marks else ""
        print(f"  {i:2}. {title_short:<55} {mark_str}")
        # 본문 일부 (긴 텍스트만)
        for t in body_texts[:2]:
            if t != title:
                print(f"      → {t}")

    return prs


# 3개 파일 분석
for pf in PPTX_FILES:
    if Path(pf).exists():
        try:
            analyze(pf)
        except Exception as e:
            print(f"❌ 분석 실패 {Path(pf).name}: {e}")
    else:
        print(f"❌ 파일 없음: {pf}")

print(f"\n{'='*70}")
print("✅ 분석 완료")
